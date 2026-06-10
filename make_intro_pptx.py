#!/usr/bin/env python3
"""Generate herbarium_pipeline_intro.pptx — run once then delete."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR

# ── Palette ────────────────────────────────────────────────────────────────
TEAL_D  = RGBColor(0x00, 0x69, 0x5C)
TEAL_M  = RGBColor(0x00, 0x89, 0x7B)
TEAL_L  = RGBColor(0x80, 0xCB, 0xC4)
TEAL_BG = RGBColor(0xE0, 0xF2, 0xF1)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x20, 0x27)
GREY    = RGBColor(0x45, 0x5A, 0x64)
LGREY   = RGBColor(0xF0, 0xF2, 0xF4)
MGREY   = RGBColor(0xCF, 0xD8, 0xDC)
ORANGE  = RGBColor(0xE6, 0x5C, 0x00)
ORANGE_L= RGBColor(0xFF, 0xE0, 0xB2)
BLUE    = RGBColor(0x01, 0x5F, 0x9B)
BLUE_L  = RGBColor(0xB3, 0xE5, 0xFC)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
GREEN_L = RGBColor(0xC8, 0xE6, 0xC9)
AMBER   = RGBColor(0xFF, 0x8F, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Low-level helpers ──────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def rrect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(1)):
    sh = slide.shapes.add_shape(5, x, y, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def box(slide, label, x, y, w, h, fill, text_color=WHITE, size=13,
        bold=False, sub=None):
    """Rounded-rect with centered label (+ optional smaller sub-label)."""
    sh = rrect(slide, x, y, w, h, fill=fill)
    sh.text_frame.word_wrap = True
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = sh.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = text_color
    if sub:
        p2 = sh.text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(size - 2); r2.font.color.rgb = text_color
    return sh

def txt(slide, text, x, y, w, h, size=13, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def bullets(slide, items, x, y, w, h, size=13, color=DARK, gap=Pt(4)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        r = p.add_run()
        prefix = "  " if item.startswith("    ") else "• "
        r.text = prefix + item.lstrip()
        r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def arrow_h(slide, x, y, length, color=MGREY, label=None):
    """Horizontal arrow using RIGHT_ARROW autoshape (id 13)."""
    h = Inches(0.28)
    sh = slide.shapes.add_shape(13, x, y - h / 2, length, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    if label:
        txt(slide, label, x, y - Inches(0.32), length, Inches(0.28),
            size=9, color=GREY, align=PP_ALIGN.CENTER)
    return sh

def arrow_v(slide, x, y, length, color=MGREY, label=None):
    """Vertical arrow using DOWN_ARROW autoshape (id 36)."""
    w = Inches(0.28)
    sh = slide.shapes.add_shape(36, x - w / 2, y, w, length)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    if label:
        txt(slide, label, x + Inches(0.06), y + length * 0.3, Inches(1.2),
            Inches(0.3), size=9, color=GREY)
    return sh

# ── Slide templates ────────────────────────────────────────────────────────

def header_slide(title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, W, Inches(1.15), fill=TEAL_D)
    rect(slide, 0, Inches(1.15), W, Inches(0.045), fill=TEAL_M)
    txt(slide, title, Inches(0.45), Inches(0.15), Inches(12.4), Inches(0.7),
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(0.78), Inches(12.4), Inches(0.35),
            size=12, color=TEAL_L)
    rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=LGREY)
    txt(slide, "Herbarium Classification Pipeline",
        Inches(0.45), H - Inches(0.36), Inches(12), Inches(0.34),
        size=10, color=GREY)
    return slide, Inches(1.28)   # returns (slide, content_top)

def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, W, H, fill=TEAL_D)
    sh = slide.shapes.add_shape(9, Inches(-1.5), Inches(4.8), Inches(4), Inches(4))
    sh.fill.solid(); sh.fill.fore_color.rgb = TEAL_M; sh.line.fill.background()
    sh = slide.shapes.add_shape(9, Inches(11.5), Inches(-1), Inches(3), Inches(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = TEAL_M; sh.line.fill.background()
    txt(slide, title, Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.4),
        size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, subtitle, Inches(1.2), Inches(3.9), Inches(10.9), Inches(0.8),
            size=17, color=TEAL_L, align=PP_ALIGN.CENTER)
    return slide

# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

# ── 1. Title ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=LGREY)
rect(slide, 0, 0, Inches(0.4), H, fill=TEAL_D)

# White centre panel
rect(slide, Inches(0.75), Inches(1.5), Inches(8.1), Inches(4.2), fill=WHITE)
rect(slide, Inches(0.75), Inches(1.5), Inches(0.12), Inches(4.2), fill=TEAL_M)

txt(slide, "Herbarium Classification\nPipeline",
    Inches(1.2), Inches(1.75), Inches(7.4), Inches(1.5),
    size=34, bold=True, color=TEAL_D)
txt(slide, "A guided walkthrough: downloading specimen images,\nfiltering, training a deep learning model, and reviewing results.",
    Inches(1.2), Inches(3.2), Inches(7.4), Inches(0.9),
    size=14, color=GREY)

stage_cols = [TEAL_D, TEAL_M, RGBColor(0x00,0xAC,0x98),
              RGBColor(0x26,0xA6,0x9A), TEAL_L]
stage_labels = ["1\nDownload", "2\nFilter", "3\nResized", "4\nTrain", "5\nIdentify"]
for i, (lbl, col) in enumerate(zip(stage_labels, stage_cols)):
    box(slide, lbl, Inches(1.2 + i * 1.38), Inches(4.85),
        Inches(1.28), Inches(0.65), fill=col, size=11)

# Right panel
rect(slide, Inches(9.2), Inches(1.0), Inches(3.8), Inches(5.7), fill=TEAL_BG)
txt(slide, "🌿", Inches(9.9), Inches(1.6), Inches(2.5), Inches(2.0),
    size=80, align=PP_ALIGN.CENTER)
bullets(slide,
    ["100s of species", "Herbarium scans", "GPU-accelerated", "Cloud or local"],
    Inches(9.45), Inches(3.8), Inches(3.3), Inches(2.5), size=13, color=GREY)

rect(slide, 0, H - Inches(0.45), W, Inches(0.45), fill=TEAL_D)
txt(slide, "George Gosline  ·  2026",
    Inches(0.5), H - Inches(0.42), Inches(12), Inches(0.38),
    size=11, color=TEAL_L)


# ── 2. What You'll Build ──────────────────────────────────────────────────
slide, ct = header_slide("What You'll Build",
    "A trained deep learning model that identifies plant species from herbarium specimen images")

box(slide, "Input\nHerbarium scans\n(any scanner format)",
    Inches(0.4), ct + Inches(0.3), Inches(3.0), Inches(2.0),
    fill=TEAL_BG, text_color=DARK, size=12)
arrow_h(slide, Inches(3.5), ct + Inches(1.3), Inches(1.3), color=TEAL_M,
        label="train model")
box(slide, "Model\ncheckpoint\n(.ckpt)",
    Inches(4.9), ct + Inches(0.3), Inches(2.6), Inches(2.0),
    fill=TEAL_M, size=12)
arrow_h(slide, Inches(7.6), ct + Inches(1.3), Inches(1.3), color=TEAL_M,
        label="inference")
box(slide, "Output\nTop-5 species\npredictions + confidence",
    Inches(9.0), ct + Inches(0.3), Inches(3.9), Inches(2.0),
    fill=TEAL_BG, text_color=DARK, size=12)

txt(slide, "Capabilities", Inches(0.5), ct + Inches(2.6), Inches(5.5), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
bullets(slide, [
    "Classifies hundreds of species simultaneously",
    "Handles hierarchical taxonomy — species / genus / family heads",
    "Optionally fuses geographic coordinates (lat/lon) with image features",
    "Flags uncertain identifications and probable mis-labelled specimens",
    "Runs on RunPod cloud GPU (~$2–4 per training run) or a local CUDA GPU",
], Inches(0.5), ct + Inches(3.0), Inches(5.8), Inches(2.8), size=13)

txt(slide, "Typical results after 15 epochs on 3 000+ specimens per species:",
    Inches(6.5), ct + Inches(2.6), Inches(6.4), Inches(0.5), size=13, color=GREY)
for i, (metric, val, col) in enumerate([
        ("Top-1 accuracy", "85–95 %", TEAL_D),
        ("Top-5 accuracy", "> 98 %", TEAL_M),
        ("Training time", "2–5 hrs", BLUE)]):
    bx = Inches(6.5 + i * 2.2)
    box(slide, f"{val}\n{metric}", bx, ct + Inches(3.1),
        Inches(2.0), Inches(1.3), fill=col, size=12)


# ── 3. Five-Stage Pipeline Overview ───────────────────────────────────────
slide, ct = header_slide("The Five-Stage Pipeline",
    "Each stage produces files consumed by the next; all persist to disk so any stage can be re-run")

stages = [
    ("1\nDownload", TEAL_D,
     "Query GBIF API\nFetch images\nWrite specsin.csv"),
    ("2\nFilter &\nCrop", RGBColor(0x00,0x6E,0x7F),
     "Remove field photos\nCLIP zero-shot classify\nCrop scanner borders"),
    ("3\nResize", RGBColor(0x01,0x57,0x9B),
     "Scale to 1 024 px\nGPU-accelerated DALI\nor PIL fallback"),
    ("4\nTrain", RGBColor(0x4A,0x14,0x8C),
     "ViT / EfficientNet\nStage 1 → Stage 2\nCheckpoints + nameslist"),
    ("5\nIdentify", RGBColor(0x1A,0x23,0x7E),
     "Run inference\npredictions.csv\nSort indets / flags"),
]

for i, (lbl, col, detail) in enumerate(stages):
    bx = Inches(0.3 + i * 2.6)
    box(slide, lbl, bx, ct + Inches(0.1), Inches(2.3), Inches(1.15),
        fill=col, size=14, bold=True)
    if i < 4:
        arrow_h(slide, bx + Inches(2.3), ct + Inches(0.65),
                Inches(0.3), color=MGREY)
    # Detail box below
    rect(slide, bx, ct + Inches(1.4), Inches(2.3), Inches(1.6),
         fill=LGREY, line_color=MGREY)
    txt(slide, detail, bx + Inches(0.1), ct + Inches(1.45),
        Inches(2.1), Inches(1.5), size=11, color=DARK)

# Data backbone row
rect(slide, Inches(0.3), ct + Inches(3.25), Inches(12.73), Inches(0.55),
     fill=TEAL_BG, line_color=TEAL_L)
txt(slide, "specsin.csv  —  the project's metadata backbone: one row per specimen, updated at each stage",
    Inches(0.5), ct + Inches(3.3), Inches(12.3), Inches(0.45),
    size=12, color=TEAL_D, align=PP_ALIGN.CENTER, bold=True)

txt(slide, "Persistent file layout under  <ProjectsRoot>/<ProjectName>/",
    Inches(0.3), ct + Inches(4.05), Inches(8), Inches(0.4),
    size=12, bold=True, color=GREY)
bullets(slide, [
    "images_cropped/   —  downloaded + filtered images",
    "runs/checkpoints/   —  model weights (.ckpt) + nameslist.json",
    "review/predictions.csv   —  inference results",
], Inches(0.3), ct + Inches(4.45), Inches(8.5), Inches(1.5), size=12, color=GREY)


# ── 4. Setup: Project Configuration ──────────────────────────────────────
slide, ct = header_slide("Step 0 — Project Setup",
    "Configure before running any pipeline stage")

# Mock header bar
rect(slide, Inches(0.4), ct + Inches(0.1), Inches(12.5), Inches(0.6),
     fill=LGREY, line_color=MGREY)
txt(slide, "Projects root:", Inches(0.6), ct + Inches(0.18), Inches(1.5), Inches(0.4),
    size=11, bold=True, color=DARK)
rect(slide, Inches(2.1), ct + Inches(0.2), Inches(3.2), Inches(0.38),
     fill=WHITE, line_color=TEAL_M)
txt(slide, "C:\\AIProjects", Inches(2.2), ct + Inches(0.22), Inches(3.0), Inches(0.34),
    size=11, color=DARK)
txt(slide, "Project name:", Inches(5.6), ct + Inches(0.18), Inches(1.3), Inches(0.4),
    size=11, bold=True, color=DARK)
rect(slide, Inches(6.9), ct + Inches(0.2), Inches(1.8), Inches(0.38),
     fill=WHITE, line_color=TEAL_M)
txt(slide, "Ebenaceae", Inches(7.0), ct + Inches(0.22), Inches(1.6), Inches(0.34),
    size=11, color=DARK)
box(slide, "Apply paths", Inches(9.4), ct + Inches(0.2), Inches(1.4), Inches(0.38),
    fill=TEAL_M, size=11)

txt(slide, "↑  The header bar is always visible — set these first, then click Apply paths.",
    Inches(0.5), ct + Inches(0.82), Inches(12), Inches(0.4),
    size=11, italic=True, color=GREY)

# Four info cards
cards = [
    ("Projects root", TEAL_BG, TEAL_D,
     "Parent folder that holds all your projects.\nExample:  C:\\AIProjects  or  D:\\herbarium"),
    ("Project name", TEAL_BG, TEAL_D,
     "Sub-folder name for this taxon run.\nExample:  Ebenaceae\nAll outputs land in  <root>/<name>/"),
    ("Image folder", BLUE_L, BLUE,
     "Which image sub-folder to use.\nDefault:  images_cropped\nChanged automatically by Apply paths."),
    ("Apply paths", GREEN_L, GREEN,
     "One click fills every stage's\ninput/output paths correctly.\nDo this first after naming the project."),
]
for i, (title, bg, accent, desc) in enumerate(cards):
    bx = Inches(0.4 + i * 3.2)
    rect(slide, bx, ct + Inches(1.3), Inches(3.0), Inches(2.5),
         fill=bg, line_color=accent, line_w=Pt(1.5))
    txt(slide, title, bx + Inches(0.15), ct + Inches(1.4), Inches(2.7), Inches(0.4),
        size=13, bold=True, color=accent)
    txt(slide, desc, bx + Inches(0.15), ct + Inches(1.8), Inches(2.7), Inches(1.5),
        size=11, color=DARK)

txt(slide, "Active model  (header card below the tabs)",
    Inches(0.5), ct + Inches(4.05), Inches(10), Inches(0.4),
    size=13, bold=True, color=TEAL_D)
bullets(slide, [
    "Paste or browse to the .ckpt checkpoint file once training is complete",
    "Used by Quick ID (drag-and-drop single-image identification) and the Identify tab",
], Inches(0.5), ct + Inches(4.45), Inches(12), Inches(1.0), size=12)


# ── 5. Stage 1: Download ──────────────────────────────────────────────────
slide, ct = header_slide("Stage 1 — Download from GBIF",
    "Queries the Global Biodiversity Information Facility API and fetches specimen images")

# Left: settings
txt(slide, "Key settings", Inches(0.4), ct + Inches(0.1), Inches(5.5), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
settings = [
    ("Taxon", "Family, Genus, or Order name\n(e.g. Ebenaceae, Diospyros, Ericales)"),
    ("Geographic filter", "Continent, country, or bounding box\nLeave blank for worldwide"),
    ("specsin.csv path", "Output metadata file — created fresh or\nappended if it already exists"),
    ("Output images dir", "Where downloaded images are saved\n(auto-filled by Apply paths)"),
    ("Limit", "Max images to download (leave blank\nfor full dataset — can be millions)"),
]
for i, (name, desc) in enumerate(settings):
    by = ct + Inches(0.55 + i * 0.92)
    rect(slide, Inches(0.4), by, Inches(5.6), Inches(0.82),
         fill=LGREY, line_color=MGREY)
    txt(slide, name, Inches(0.55), by + Inches(0.06), Inches(1.8), Inches(0.34),
        size=12, bold=True, color=TEAL_D)
    txt(slide, desc, Inches(2.4), by + Inches(0.06), Inches(3.5), Inches(0.7),
        size=11, color=DARK)

# Right: output diagram
txt(slide, "What Download produces", Inches(6.5), ct + Inches(0.1),
    Inches(6.5), Inches(0.4), size=14, bold=True, color=TEAL_D)
box(slide, "GBIF API\ngbif.org", Inches(6.6), ct + Inches(0.55),
    Inches(2.2), Inches(1.0), fill=GREEN, size=12)
arrow_v(slide, Inches(7.7), ct + Inches(1.55), Inches(0.6), color=GREEN,
        label="HTTP")
box(slide, "specsin.csv\n(one row per specimen)",
    Inches(6.6), ct + Inches(2.2), Inches(2.2), Inches(1.0),
    fill=TEAL_M, size=11)
box(slide, "images_cropped/\n(JPEGs on disk)",
    Inches(9.2), ct + Inches(2.2), Inches(2.8), Inches(1.0),
    fill=TEAL_D, size=11)

txt(slide, "specsin.csv columns added at download:",
    Inches(6.5), ct + Inches(3.5), Inches(6.5), Inches(0.38), size=12,
    bold=True, color=GREY)
bullets(slide, [
    "catalogNumber, species, family, genus, verbatimName",
    "decimalLatitude, decimalLongitude, countryCode",
    "gbifID, image_url, fname, hasfile",
], Inches(6.5), ct + Inches(3.9), Inches(6.5), Inches(1.5), size=12, color=GREY)


# ── 6. Stage 2: Filter & Crop ─────────────────────────────────────────────
slide, ct = header_slide("Stage 2 — Filter & Crop",
    "Removes non-herbarium images and crops dark scanner borders")

col1 = Inches(0.4)
col2 = Inches(6.8)

txt(slide, "What gets removed", col1, ct + Inches(0.05), Inches(5.8), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
remove_items = [
    ("🌿 Field photos", "Living-plant photographs taken in the field"),
    ("🔬 Slides / close-ups", "Microscope slides and detail scans"),
    ("📄 Label-only scans", "Pages with no specimen material"),
    ("🖼 Low-quality images", "Very dark, blurred, or tiny thumbnails"),
]
for i, (icon_lbl, desc) in enumerate(remove_items):
    by = ct + Inches(0.5 + i * 0.8)
    box(slide, icon_lbl, col1, by, Inches(2.0), Inches(0.6),
        fill=RGBColor(0xFF,0xEB,0xEE), text_color=RGBColor(0xC6,0x28,0x28), size=12)
    txt(slide, desc, col1 + Inches(2.15), by + Inches(0.12),
        Inches(3.5), Inches(0.4), size=12, color=DARK)

txt(slide, "Two filter methods", col1, ct + Inches(3.35), Inches(5.8), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
box(slide, "CLIP  (default)\nZero-shot image classification\nusing a vision-language model\n\n✓ Best accuracy\n✗ Requires GPU",
    col1, ct + Inches(3.75), Inches(2.7), Inches(2.0),
    fill=TEAL_BG, text_color=DARK, size=11)
box(slide, "HSV heuristic\nColour-space analysis\nfor sheet-white backgrounds\n\n✓ CPU-only, fast\n✗ Misses coloured mounts",
    col1 + Inches(2.9), ct + Inches(3.75), Inches(2.7), Inches(2.0),
    fill=LGREY, text_color=DARK, size=11)

# Right side: crop diagram
txt(slide, "Border cropping", col2, ct + Inches(0.05), Inches(6.1), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
# Before box
rect(slide, col2, ct + Inches(0.5), Inches(2.8), Inches(2.8),
     fill=DARK, line_color=GREY)
rect(slide, col2 + Inches(0.4), ct + Inches(0.9), Inches(2.0), Inches(2.0),
     fill=WHITE, line_color=MGREY)
txt(slide, "Before\n(dark scanner border)",
    col2, ct + Inches(3.35), Inches(2.8), Inches(0.55),
    size=11, color=GREY, align=PP_ALIGN.CENTER)
arrow_h(slide, col2 + Inches(2.95), ct + Inches(1.9), Inches(0.9),
        color=TEAL_M)
# After box
rect(slide, col2 + Inches(4.0), ct + Inches(0.9), Inches(2.0), Inches(2.0),
     fill=WHITE, line_color=TEAL_M)
txt(slide, "After\n(borders removed)",
    col2 + Inches(3.75), ct + Inches(3.35), Inches(2.8), Inches(0.55),
    size=11, color=GREY, align=PP_ALIGN.CENTER)

txt(slide, "Rejected images are moved to  images_cropped/rejected/  (field photos → live/)\nnot deleted — you can inspect and recover them.",
    col2, ct + Inches(4.15), Inches(6.3), Inches(0.8),
    size=12, italic=True, color=GREY)


# ── 7. Stage 3: Resize ────────────────────────────────────────────────────
slide, ct = header_slide("Stage 3 — Resize",
    "Standardises image dimensions so the training data-loader can form uniform batches")

bullets(slide, [
    "Scales every image so its longest edge ≤ 1 024 px  (configurable via --max-size)",
    "Maintains aspect ratio — no distortion",
    "Uses NVIDIA DALI GPU pipeline when available; falls back to PIL on CPU",
    "Overwrites images in-place (--in-place) or writes to a separate output folder",
    "Images already smaller than the target are left unchanged",
], Inches(0.5), ct + Inches(0.1), Inches(7.8), Inches(2.8), size=14)

# Size comparison diagram
sizes = [("Raw scan\n~4 000 px", Inches(2.5), ORANGE),
         ("Resized\n1 024 px", Inches(1.0), TEAL_M),
         ("Model input\n640 px  (--image-sz)", Inches(0.65), BLUE)]
base_y = ct + Inches(3.2)
bx = Inches(1.0)
for i, (lbl, sz, col) in enumerate(sizes):
    rect(slide, bx, base_y + (Inches(1.5) - sz) / 2,
         sz, sz, fill=col, line_color=WHITE)
    txt(slide, lbl, bx, base_y + Inches(1.6),
        sz + Inches(0.3), Inches(0.5), size=11, color=DARK,
        align=PP_ALIGN.CENTER)
    if i < 2:
        arrow_h(slide, bx + sz + Inches(0.1),
                base_y + Inches(0.75), Inches(0.5), color=MGREY)
    bx += sz + Inches(0.7)

txt(slide, "Why 1 024 px?",
    Inches(6.5), ct + Inches(0.1), Inches(6.5), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
bullets(slide, [
    "Herbarium scans carry fine detail (venation, surface texture)",
    "ViT-Large with --image-sz 640 still sees far more detail than most photo datasets",
    "GPU VRAM is the bottleneck — 1 024 px source with 640 px training crop is the sweet spot",
    "Batch size 4–8 fits in 20 GB VRAM at 640 px; reduce to 2–4 at 768 px",
], Inches(6.5), ct + Inches(0.55), Inches(6.5), Inches(2.5), size=13)

txt(slide, "Tip  — you can skip this step entirely if images are already ≤ 1 024 px or if VRAM allows full-res.",
    Inches(0.5), H - Inches(1.2), Inches(12.3), Inches(0.5),
    size=12, italic=True, color=GREY)


# ── 8. Stage 4: Train ─────────────────────────────────────────────────────
slide, ct = header_slide("Stage 4 — Train the Model",
    "Multi-stage fine-tuning of a pretrained vision model on your labelled herbarium images")

# Training stages diagram
for i, (lbl, sub, col) in enumerate([
    ("Stage 1\nWarm-up", "Backbone frozen\nHead only  ·  4 epochs\nLR ~ 0.005", TEAL_M),
    ("Stage 2\nFine-tune", "All layers train\n15–50 epochs\nLR ~ 0.0001", TEAL_D),
    ("Cool-down\n(optional)", "Reduced batch + LR\nFinal polishing\nSaves best-cd.ckpt", BLUE),
]):
    bx = Inches(0.4 + i * 4.2)
    box(slide, lbl, bx, ct + Inches(0.05), Inches(3.8), Inches(0.75),
        fill=col, size=13, bold=True)
    rect(slide, bx, ct + Inches(0.85), Inches(3.8), Inches(1.1),
         fill=LGREY, line_color=MGREY)
    txt(slide, sub, bx + Inches(0.1), ct + Inches(0.9), Inches(3.6), Inches(1.0),
        size=12, color=DARK)
    if i < 2:
        arrow_h(slide, bx + Inches(3.8), ct + Inches(0.42),
                Inches(0.4), color=TEAL_L)

# Key hyperparameter cards
txt(slide, "Key hyperparameters", Inches(0.4), ct + Inches(2.2),
    Inches(12), Inches(0.4), size=14, bold=True, color=TEAL_D)
params = [
    ("Model", "vit_large_patch16_dinov2\n(default)\nor efficientnet_b4 for speed"),
    ("Image size\n--image-sz", "640 px  (default)\n768 for highest accuracy\n512 for low VRAM"),
    ("Batch size", "4–8 on RTX 4090\nIncrease --grad-accum\nto compensate if reduced"),
    ("Stage 2 epochs", "15 (quick test)\n30 for production\n50 for large datasets"),
]
for i, (name, val) in enumerate(params):
    bx = Inches(0.4 + i * 3.2)
    rect(slide, bx, ct + Inches(2.65), Inches(3.0), Inches(2.1),
         fill=TEAL_BG, line_color=TEAL_L)
    txt(slide, name, bx + Inches(0.1), ct + Inches(2.72), Inches(2.8), Inches(0.55),
        size=12, bold=True, color=TEAL_D)
    txt(slide, val, bx + Inches(0.1), ct + Inches(3.2), Inches(2.8), Inches(1.2),
        size=12, color=DARK)

txt(slide, "Outputs:  runs/checkpoints/last.ckpt  ·  runs/checkpoints/epoch=XX-valid_loss=X.ckpt  ·  runs/nameslist.json",
    Inches(0.4), H - Inches(1.0), Inches(12.5), Inches(0.5),
    size=11, color=GREY, italic=True)


# ── 9. Stage 5: Identify & Review ─────────────────────────────────────────
slide, ct = header_slide("Stage 5 — Identify & Review",
    "Run inference on your images, inspect results, and correct any mislabelled specimens")

txt(slide, "Identify tab", Inches(0.4), ct + Inches(0.05), Inches(6.0), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
bullets(slide, [
    "Point to your checkpoint (.ckpt) and specsin.csv",
    "Outputs  review/predictions.csv  with top-5 predictions + confidence per image",
    "Flags images where the top prediction differs from the recorded label  (mismatch column)",
    "Sorts unidentified specimens (indet) by prediction confidence into  review/indets/",
], Inches(0.4), ct + Inches(0.48), Inches(6.0), Inches(2.2), size=13)

txt(slide, "Review tab", Inches(6.8), ct + Inches(0.05), Inches(6.2), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
bullets(slide, [
    "Browse images with their predictions side-by-side",
    "Click a prediction to confirm or correct the label",
    "Changes saved back to  predictions.csv  in real time",
], Inches(6.8), ct + Inches(0.48), Inches(6.2), Inches(1.5), size=13)

txt(slide, "Analysis tab", Inches(6.8), ct + Inches(2.1), Inches(6.2), Inches(0.4),
    size=14, bold=True, color=TEAL_D)
bullets(slide, [
    "Load  predictions.csv  to plot confusion matrix",
    "Per-species accuracy, top confusions, recall vs. support",
], Inches(6.8), ct + Inches(2.53), Inches(6.2), Inches(1.2), size=13)

# predictions.csv structure
txt(slide, "predictions.csv columns", Inches(0.4), ct + Inches(2.9),
    Inches(6.0), Inches(0.4), size=14, bold=True, color=TEAL_D)
cols = ["path", "true_label", "top_pred", "top_conf",
        "pred_2 … pred_5", "conf_2 … conf_5", "flagged", "mismatch"]
for i, c in enumerate(cols):
    bx = Inches(0.4 + (i % 4) * 1.52)
    by = ct + Inches(3.35 + (i // 4) * 0.5)
    box(slide, c, bx, by, Inches(1.42), Inches(0.4),
        fill=TEAL_BG if i < 3 else LGREY,
        text_color=TEAL_D if i < 3 else GREY, size=10)


# ── 10. Run All ───────────────────────────────────────────────────────────
slide, ct = header_slide("Run All — Full Pipeline Automation",
    "Executes all five stages in sequence; can be paused, cancelled, and resumed")

for i, (lbl, col) in enumerate(zip(
        ["1 Download", "2 Filter & Crop", "3 Resize", "4 Train", "5 Identify"],
        [TEAL_D, RGBColor(0x00,0x6E,0x7F), BLUE, RGBColor(0x4A,0x14,0x8C),
         RGBColor(0x1A,0x23,0x7E)])):
    bx = Inches(0.4 + i * 2.55)
    box(slide, lbl, bx, ct + Inches(0.2), Inches(2.35), Inches(0.7),
        fill=col, size=13, bold=True)
    if i < 4:
        arrow_h(slide, bx + Inches(2.35), ct + Inches(0.55),
                Inches(0.2), color=MGREY)

bullets(slide, [
    "Each stage checks whether its output already exists — skips completed stages automatically",
    "If a stage fails, fix the issue and restart Run All; completed stages are not repeated",
    "The Stop button cancels the current stage cleanly (no partial-write corruption)",
    "In Cloud mode the orchestrator runs all stages on the pod; progress streams to the Output panel",
], Inches(0.4), ct + Inches(1.15), Inches(12.5), Inches(2.0), size=13)

# Mode comparison mini-table
for col_x, mode, bg, steps in [
        (Inches(0.4), "Local mode", TEAL_BG,
         ["Scripts run on your own machine",
          "Needs CUDA GPU for training",
          "No account setup beyond GBIF",
          "Output stays in your project folder"]),
        (Inches(6.8), "Cloud mode  (☁)", BLUE_L,
         ["Orchestrator provisions a RunPod pod",
          "Streams output to the web UI log panel",
          "Archives results to Cloudflare R2",
          "Pod can be terminated when done"])]:
    rect(slide, col_x, ct + Inches(3.3), Inches(6.0), Inches(2.4),
         fill=bg, line_color=TEAL_M if "Local" in mode else BLUE)
    txt(slide, mode, col_x + Inches(0.15), ct + Inches(3.38), Inches(5.7), Inches(0.4),
        size=13, bold=True, color=TEAL_D if "Local" in mode else BLUE)
    bullets(slide, steps, col_x + Inches(0.1), ct + Inches(3.78),
            Inches(5.7), Inches(1.7), size=12,
            color=TEAL_D if "Local" in mode else BLUE)


# ── Section: Cloud Mode ───────────────────────────────────────────────────
section_slide("Cloud Mode",
    "GPU training on RunPod · Storage on Cloudflare R2 · Orchestrated from your browser")


# ── 11. Cloud vs Local ────────────────────────────────────────────────────
slide, ct = header_slide("Cloud vs Local — Choosing a Mode")

rows = [
    ("", "Cloud mode  ☁", "Local mode  💻"),
    ("GPU", "RunPod RTX 4090 (provisioned on demand)", "Your own NVIDIA GPU (≥ 20 GB VRAM recommended)"),
    ("Cost", "~$2–4 per training run  (pod time + R2 storage)", "Electricity only"),
    ("Setup", "RunPod + R2 accounts, SSH key, API keys", "conda env + CUDA drivers"),
    ("Data storage", "Network volume on pod + Cloudflare R2 archive", "Local disk"),
    ("Monitoring", "Live log stream in browser from any device", "Terminal / Output panel on same machine"),
    ("Resuming", "Archive to R2 → restore later with restore_local.py", "Checkpoint on disk"),
    ("Best for", "Long runs, no local GPU, remote collaboration", "Quick iteration, offline, sensitive data"),
]
col_ws = [Inches(2.4), Inches(5.0), Inches(5.0)]
col_xs = [Inches(0.4), Inches(2.85), Inches(7.9)]
for r_i, row in enumerate(rows):
    bg = TEAL_D if r_i == 0 else (TEAL_BG if r_i % 2 == 1 else WHITE)
    tc = WHITE if r_i == 0 else DARK
    by = ct + Inches(0.05 + r_i * 0.64)
    for c_i, (cell, cw, cx) in enumerate(zip(row, col_ws, col_xs)):
        rect(slide, cx, by, cw, Inches(0.62),
             fill=bg, line_color=MGREY, line_w=Pt(0.5))
        txt(slide, cell, cx + Inches(0.08), by + Inches(0.08),
            cw - Inches(0.12), Inches(0.5),
            size=12 if r_i > 0 else 13,
            bold=(r_i == 0 or c_i == 0),
            color=tc)


# ── 12. Cloud Architecture ────────────────────────────────────────────────
slide, ct = header_slide("Cloud Architecture",
    "Three main components communicate over the internet; your browser is the control plane")

by = ct + Inches(0.3)
bh = Inches(3.2)

# Your PC
rect(slide, Inches(0.2), by, Inches(3.3), bh,
     fill=TEAL_BG, line_color=TEAL_D, line_w=Pt(2))
txt(slide, "Your PC", Inches(0.3), by + Inches(0.08), Inches(3.1), Inches(0.45),
    size=14, bold=True, color=TEAL_D)
bullets(slide, ["Web UI (browser)", "Orchestrator process", "rclone (restore)", "SSH key"],
        Inches(0.35), by + Inches(0.55), Inches(2.9), Inches(2.2), size=12)

# Arrows PC ↔ Pod
arrow_h(slide, Inches(3.55), by + Inches(1.2),
        Inches(1.4), color=TEAL_M, label="SSH / SFTP")
arrow_h(slide, Inches(3.55), by + Inches(1.65),
        Inches(1.4), color=TEAL_L, label="log stream")

# RunPod Pod
rect(slide, Inches(5.1), by, Inches(3.9), bh,
     fill=BLUE_L, line_color=BLUE, line_w=Pt(2))
txt(slide, "RunPod GPU Pod", Inches(5.2), by + Inches(0.08), Inches(3.7), Inches(0.45),
    size=14, bold=True, color=BLUE)
bullets(slide, ["RTX 4090 / A100", "Download  ·  Filter  ·  Resize",
                "Train  ·  Identify", "rclone (archive to R2)"],
        Inches(5.25), by + Inches(0.55), Inches(3.5), Inches(2.2), size=12)

# Arrows Pod ↔ R2
arrow_h(slide, Inches(9.1), by + Inches(1.2),
        Inches(1.3), color=ORANGE, label="rclone copy")
arrow_h(slide, Inches(9.1), by + Inches(1.65),
        Inches(1.3), color=RGBColor(0xFF,0xCC,0x80), label="rclone restore")

# R2
rect(slide, Inches(10.5), by, Inches(2.6), bh,
     fill=ORANGE_L, line_color=ORANGE, line_w=Pt(2))
txt(slide, "Cloudflare R2", Inches(10.6), by + Inches(0.08), Inches(2.4), Inches(0.45),
    size=14, bold=True, color=ORANGE)
bullets(slide, ["checkpoints/", "images.tar", "specsin.csv", "nameslist.json"],
        Inches(10.65), by + Inches(0.55), Inches(2.3), Inches(2.2), size=12)

# GBIF below pod
arrow_v(slide, Inches(7.05), by + bh, Inches(0.5),
        color=GREEN, label="HTTP")
box(slide, "GBIF API\n(specimen records + image URLs)",
    Inches(5.1), by + bh + Inches(0.55), Inches(3.9), Inches(0.8),
    fill=GREEN_L, text_color=GREEN, size=12)

txt(slide, "Note: the orchestrator runs on Your PC and issues commands to the pod over SSH."
    "  The web UI is also served locally — your browser connects to localhost:8765.",
    Inches(0.3), H - Inches(1.0), Inches(12.7), Inches(0.55),
    size=11, italic=True, color=GREY)


# ── 13. Data Flow: Download Phase ─────────────────────────────────────────
slide, ct = header_slide("Data Flow — Download & Prepare Phases",
    "GBIF → Pod disk  →  R2 archive  (stages 1 – 3)")

steps = [
    ("1  Provision", "Orchestrator calls RunPod API.\nPod boots, SSH comes up (~2 min).\nrclone config pushed to pod.", TEAL_D),
    ("2  Download", "Pod queries GBIF API.\nImages fetched in parallel threads.\nspecsin.csv built incrementally.", GREEN),
    ("3  Filter & Crop", "CLIP runs on pod GPU.\nNon-herbarium images moved to rejected/.\nBorders cropped.", BLUE),
    ("4  Resize", "DALI GPU pipeline scales images.\nImages overwritten in-place.", RGBColor(0x4A,0x14,0x8C)),
    ("5  Archive", "rclone copies images.tar + specsin.csv\nto R2.  Pod can now be\nterminated safely.", ORANGE),
]
for i, (lbl, detail, col) in enumerate(steps):
    bx = Inches(0.25 + i * 2.6)
    box(slide, lbl, bx, ct + Inches(0.05), Inches(2.4), Inches(0.6),
        fill=col, size=12, bold=True)
    if i < 4:
        arrow_h(slide, bx + Inches(2.4), ct + Inches(0.35),
                Inches(0.2), color=MGREY)
    rect(slide, bx, ct + Inches(0.72), Inches(2.4), Inches(1.8),
         fill=LGREY, line_color=MGREY)
    txt(slide, detail, bx + Inches(0.1), ct + Inches(0.78),
        Inches(2.2), Inches(1.7), size=11, color=DARK)

# Storage diagram
txt(slide, "Where data lives after each phase",
    Inches(0.3), ct + Inches(2.85), Inches(12.7), Inches(0.38),
    size=13, bold=True, color=TEAL_D)
for ix, (lbl, sub, col) in enumerate([
    ("Pod network\nvolume\n/workspace/data/", "images_cropped/\nspecsin.csv\n(temporary)", BLUE),
    ("Pod ephemeral\n/workspace/data/\n(deleted on terminate)", "images_1024/\nresized images", RGBColor(0x4A,0x14,0x8C)),
    ("Cloudflare R2\nherbarium-backup/\n<project>/", "images.tar\nspecsin.csv\n(permanent)", ORANGE),
]):
    bx = Inches(0.4 + ix * 4.2)
    rect(slide, bx, ct + Inches(3.3), Inches(3.9), Inches(2.1),
         fill=BLUE_L if col == BLUE else (ORANGE_L if col == ORANGE else LGREY),
         line_color=col, line_w=Pt(1.5))
    txt(slide, lbl, bx + Inches(0.12), ct + Inches(3.38), Inches(3.6), Inches(0.9),
        size=12, bold=True, color=col)
    txt(slide, sub, bx + Inches(0.12), ct + Inches(4.22), Inches(3.6), Inches(0.9),
        size=11, color=DARK, italic=True)


# ── 14. Data Flow: Training & Archive ─────────────────────────────────────
slide, ct = header_slide("Data Flow — Training & Archive",
    "Stage 4: pod GPU reads images, writes checkpoints; R2 stores the final model")

# Training loop diagram
pod_x = Inches(1.5); pod_y = ct + Inches(0.2); pod_w = Inches(5.0); pod_h = Inches(3.0)
rect(slide, pod_x, pod_y, pod_w, pod_h, fill=BLUE_L, line_color=BLUE, line_w=Pt(2))
txt(slide, "RunPod GPU Pod  —  Training Loop",
    pod_x + Inches(0.15), pod_y + Inches(0.1), Inches(4.7), Inches(0.4),
    size=13, bold=True, color=BLUE)
box(slide, "DataLoader\nreads batches\nfrom disk",
    pod_x + Inches(0.15), pod_y + Inches(0.6), Inches(1.8), Inches(1.0),
    fill=BLUE, size=11)
arrow_h(slide, pod_x + Inches(2.05), pod_y + Inches(1.1),
        Inches(0.6), color=BLUE)
box(slide, "ViT-Large\nforward + backward\nGPU compute",
    pod_x + Inches(2.7), pod_y + Inches(0.6), Inches(2.0), Inches(1.0),
    fill=RGBColor(0x4A,0x14,0x8C), size=11)
arrow_v(slide, pod_x + Inches(3.7), pod_y + Inches(1.6),
        Inches(0.8), color=RGBColor(0x4A,0x14,0x8C))
box(slide, "Checkpoint\n.ckpt saved\nper epoch",
    pod_x + Inches(2.7), pod_y + Inches(2.45), Inches(2.0), Inches(0.5),
    fill=TEAL_D, size=11)

# Arrow pod → R2
arrow_h(slide, pod_x + pod_w + Inches(0.1), pod_y + Inches(1.5),
        Inches(1.3), color=ORANGE, label="rclone\n(best ckpt)")

# R2 box
rect(slide, Inches(8.1), pod_y, Inches(2.5), pod_h,
     fill=ORANGE_L, line_color=ORANGE, line_w=Pt(2))
txt(slide, "R2\nherbarium-backup/", Inches(8.2), pod_y + Inches(0.1),
    Inches(2.3), Inches(0.55), size=12, bold=True, color=ORANGE)
bullets(slide, ["last.ckpt", "best epoch ckpt", "nameslist.json", "specsin.csv"],
        Inches(8.2), pod_y + Inches(0.7), Inches(2.2), Inches(2.0), size=11)

# Arrow R2 → local
arrow_h(slide, Inches(10.7), pod_y + Inches(1.5),
        Inches(1.5), color=TEAL_M, label="restore_local.py")
box(slide, "Your PC\n(for Review)",
    Inches(12.3), pod_y + Inches(0.9), Inches(0.8), Inches(1.2),
    fill=TEAL_M, size=10)

# WandB note
rect(slide, Inches(1.5), ct + Inches(3.45), Inches(5.0), Inches(1.1),
     fill=LGREY, line_color=AMBER, line_w=Pt(1.5))
txt(slide, "Optional: WandB live logging",
    Inches(1.65), ct + Inches(3.52), Inches(4.7), Inches(0.38),
    size=12, bold=True, color=AMBER)
txt(slide, "Metrics stream from pod → WandB cloud → your browser\n"
    "Set WandB API key in the Setup tab to enable",
    Inches(1.65), ct + Inches(3.88), Inches(4.7), Inches(0.6),
    size=11, color=DARK)

bullets(slide, [
    "Stage 1  (4 epochs, head only)  →  Stage 2  (15–50 epochs, full fine-tune)",
    "Best-epoch checkpoint is identified by valid_loss and archived to R2 after Stage 2",
    "Resume from checkpoint: if  last.ckpt  exists, Stage 1 is skipped automatically",
], Inches(0.4), ct + Inches(4.75), Inches(12.5), Inches(1.5), size=12, color=GREY)


# ── 15. Data Flow: Restore & Review ───────────────────────────────────────
slide, ct = header_slide("Data Flow — Restore & Local Review",
    "Pull a completed project from R2 to your PC for inspection and correction")

steps_r = [
    ("Run\nrestore_local.py", "Reads R2 credentials\nfrom OS keyring.\nNo rclone.conf needed.", TEAL_D),
    ("Download\ncheckpoints/", "last.ckpt + nameslist.json\npulled from\nR2 → local disk.", BLUE),
    ("Download\nspecsin.csv", "Project metadata\nrestored to\ntarget folder.", TEAL_M),
    ("Extract\nimages.tar", "Full image set\nextracted to\nimages/ folder.", GREEN),
    ("Open\nweb UI", "Point Active model\nto the .ckpt, click\nApply paths → Identify.", ORANGE),
]
for i, (lbl, detail, col) in enumerate(steps_r):
    bx = Inches(0.25 + i * 2.6)
    box(slide, lbl, bx, ct + Inches(0.05), Inches(2.4), Inches(0.65),
        fill=col, size=12, bold=True)
    if i < 4:
        arrow_h(slide, bx + Inches(2.4), ct + Inches(0.37),
                Inches(0.2), color=MGREY)
    rect(slide, bx, ct + Inches(0.77), Inches(2.4), Inches(1.6),
         fill=LGREY, line_color=MGREY)
    txt(slide, detail, bx + Inches(0.1), ct + Inches(0.82),
        Inches(2.2), Inches(1.5), size=11, color=DARK)

txt(slide, "Restore command", Inches(0.4), ct + Inches(2.65),
    Inches(12), Inches(0.4), size=13, bold=True, color=TEAL_D)
rect(slide, Inches(0.4), ct + Inches(3.05), Inches(12.2), Inches(0.65),
     fill=DARK)
txt(slide, 'python restore_local.py  --project Ebenaceae  --target C:\\AIProjects\\Ebenaceae  --remote r2:herbarium-backup',
    Inches(0.6), ct + Inches(3.1), Inches(12.0), Inches(0.55),
    size=12, color=RGBColor(0x80,0xCB,0xC4))

txt(slide, "After restore the target folder contains:",
    Inches(0.4), ct + Inches(3.95), Inches(5.5), Inches(0.38),
    size=13, bold=True, color=GREY)
bullets(slide, [
    "checkpoints/last.ckpt  +  nameslist.json",
    "specsin.csv",
    "images/   (full image set)",
    "predictions/   (if already run Identify)",
], Inches(0.4), ct + Inches(4.35), Inches(5.5), Inches(1.6), size=12, color=GREY)

txt(slide, "Then in the web UI:",
    Inches(6.5), ct + Inches(3.95), Inches(6.2), Inches(0.38),
    size=13, bold=True, color=TEAL_D)
bullets(slide, [
    "Header → Active model → browse to  last.ckpt",
    "Header → Projects root + Project name → Apply paths",
    "Identify tab → Run  (generates predictions.csv)",
    "Review tab → inspect predictions, correct labels",
    "Analysis tab → confusion matrix + per-species accuracy",
], Inches(6.5), ct + Inches(4.35), Inches(6.2), Inches(1.8), size=12)


# ── 16. Getting Started Checklist ─────────────────────────────────────────
slide, ct = header_slide("Getting Started — Quick-Start Checklist")

txt(slide, "Accounts & credentials", Inches(0.4), ct + Inches(0.05),
    Inches(6.0), Inches(0.4), size=14, bold=True, color=TEAL_D)
acct_items = [
    (ORANGE, "Cloudflare R2", "Sign up · create herbarium-backup bucket · generate API token (Read & Edit)"),
    (BLUE,   "RunPod",        "Sign up · add billing · create API key (rpa_…) · register SSH public key"),
    (AMBER,  "WandB",         "Optional — free for academic use · get API key for live training graphs"),
]
for i, (col, svc, desc) in enumerate(acct_items):
    by = ct + Inches(0.5 + i * 0.75)
    box(slide, svc, Inches(0.4), by, Inches(1.6), Inches(0.58), fill=col, size=12)
    txt(slide, desc, Inches(2.15), by + Inches(0.1), Inches(4.2), Inches(0.45),
        size=12, color=DARK)

txt(slide, "One-time local setup", Inches(0.4), ct + Inches(2.8),
    Inches(6.0), Inches(0.4), size=14, bold=True, color=TEAL_D)
bullets(slide, [
    "Generate SSH key:   ssh-keygen -t ed25519 -f %USERPROFILE%\\.ssh\\id_ed25519_herbarium -N \"\"",
    "Register the .pub content at RunPod → Settings → SSH Public Keys",
    "Enter keys in the web UI ⚙ Setup tab and click Save for each section",
    "Run  python herbarium_pipeline_webui.py  — opens at http://localhost:8765",
], Inches(0.4), ct + Inches(3.2), Inches(6.2), Inches(2.0), size=12)

txt(slide, "First run", Inches(6.8), ct + Inches(0.05),
    Inches(6.2), Inches(0.4), size=14, bold=True, color=TEAL_D)
steps_f = [
    "Set Projects root and Project name in the header",
    "Click Apply paths",
    "Download tab → choose a Family (start small: limit 200)",
    "Filter & Crop tab → Run",
    "Resize tab → Run",
    "Train tab → set Stage 2 epochs to 5 for a test run",
    "Identify tab → Run and check predictions.csv",
]
for i, s in enumerate(steps_f):
    by = ct + Inches(0.5 + i * 0.62)
    box(slide, str(i + 1), Inches(6.8), by, Inches(0.42), Inches(0.52),
        fill=TEAL_D, size=13, bold=True)
    txt(slide, s, Inches(7.35), by + Inches(0.08), Inches(5.7), Inches(0.44),
        size=12, color=DARK)


# ── Save ──────────────────────────────────────────────────────────────────
out = "herbarium_pipeline_intro.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
